#!/usr/bin/env python3
"""
Multi-Agent IT Assistant — Agenticthon Proposal Pitch Deck Generator
==========================================================
Generates docs/pitch-deck.pptx via python-pptx. Reproducible — run from
the repo root:

    python3 scripts/build_pitch_deck.py

Design system, copy, and slide structure are the proposal-phase decisions
locked in this session. Do not edit the .pptx by hand; edit this script
and re-run.
"""

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor


# ── Design system ────────────────────────────────────────────────────────

BG              = RGBColor(0x0A, 0x0A, 0x0A)
CARD_BG         = RGBColor(0x12, 0x12, 0x12)
TEXT_PRIMARY    = RGBColor(0xF5, 0xF5, 0xF5)
TEXT_SECONDARY  = RGBColor(0x9C, 0xA3, 0xAF)
DIVIDER         = RGBColor(0x1F, 0x1F, 0x1F)

MEMORY   = RGBColor(0x4A, 0x9E, 0xFF)  # blue
RESOLVER = RGBColor(0x4A, 0xDE, 0x80)  # green
GUARDIAN = RGBColor(0xF8, 0x71, 0x71)  # red
REPORTER = RGBColor(0xA7, 0x8B, 0xFA)  # purple
AMBER    = RGBColor(0xFB, 0xBF, 0x24)  # in-progress phase

# Single-family Arabic stack for a cleaner, more professional read.
# IBM Plex Sans Arabic is widely available and renders crisply at small
# sizes; falls back gracefully on machines without it.
FONT_HEAD = "IBM Plex Sans Arabic"
FONT_BODY = "IBM Plex Sans Arabic"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

OUT = Path(__file__).resolve().parents[1] / "docs" / "pitch-deck.pptx"


# ── Helpers ──────────────────────────────────────────────────────────────


def set_rtl(paragraph):
    """Mark a paragraph as right-to-left so PowerPoint shapes Arabic correctly."""
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set("rtl", "1")


def add_solid_bg(slide, color=BG):
    """Fill the slide with a solid color and send to back."""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    sp = bg._element
    tree = sp.getparent()
    tree.remove(sp)
    tree.insert(2, sp)
    return bg


def add_text(slide, x, y, w, h, *runs, align=PP_ALIGN.RIGHT, rtl=True,
             anchor=MSO_ANCHOR.TOP):
    """
    Add a text box. Each item in *runs is a dict describing one run:

        {"text": str,
         "font": "Cairo" | "Tajawal" | None  (defaults to Tajawal),
         "size": int (pt),
         "color": RGBColor,
         "bold": bool, "italic": bool,
         "newline": bool — start a new paragraph BEFORE this run}

    The first run goes into the existing first paragraph; subsequent runs
    with `newline=True` start a fresh paragraph (carrying the same
    alignment + RTL flag).
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    if rtl:
        set_rtl(p)

    for spec in runs:
        if spec.get("newline"):
            p = tf.add_paragraph()
            p.alignment = align
            if rtl:
                set_rtl(p)
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", FONT_BODY)
        run.font.size = Pt(spec.get("size", 18))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = spec.get("color", TEXT_PRIMARY)
    return tb


def add_dot(slide, cx, cy, radius, color):
    """Solid-color circle (for the four-agent decoration)."""
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                               cx - radius, cy - radius,
                               radius * 2, radius * 2)
    s.line.fill.background()
    s.fill.solid()
    s.fill.fore_color.rgb = color
    return s


def add_hline(slide, x1, y, x2, *, color=DIVIDER, weight=0.75):
    line = slide.shapes.add_connector(1, x1, y, x2, y)  # 1 = STRAIGHT
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_card(slide, x, y, w, h, *, accent=None, fill=CARD_BG):
    """Rounded-rect card. If `accent` is set, draws a thin colored bar
    along the top edge."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.adjustments[0] = 0.04
    card.line.color.rgb = DIVIDER
    card.line.width = Pt(0.75)
    card.fill.solid()
    card.fill.fore_color.rgb = fill
    if accent is not None:
        bar_h = Inches(0.06)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, bar_h)
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
    return card


def add_slide_title(slide, ar, en):
    """Stock top-right title block used on slides 2-7."""
    add_text(slide, Inches(7.5), Inches(0.45), Inches(5.4), Inches(0.85),
             {"text": ar, "font": FONT_HEAD, "size": 44,
              "bold": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.RIGHT, rtl=True)
    add_text(slide, Inches(7.5), Inches(1.30), Inches(5.4), Inches(0.4),
             {"text": en, "font": FONT_BODY, "size": 18,
              "color": TEXT_SECONDARY},
             align=PP_ALIGN.RIGHT, rtl=False)


# ── Slide builders ───────────────────────────────────────────────────────


def slide_title(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    cw = SLIDE_W

    # Top tagline
    add_text(s, Inches(2), Inches(0.55), cw - Inches(4), Inches(0.4),
             {"text": "Agenticthon — Multi-Agent Systems Track",
              "font": FONT_BODY, "size": 14, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)

    # Heading (very large Arabic)
    add_text(s, Inches(0.5), Inches(1.7), cw - Inches(1), Inches(1.4),
             {"text": "المساعد التقني متعدد الوكلاء",
              "font": FONT_HEAD, "size": 56,
              "bold": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)

    # English subtitle
    add_text(s, Inches(2), Inches(3.1), cw - Inches(4), Inches(0.5),
             {"text": "Multi-Agent IT Assistant",
              "font": FONT_BODY, "size": 24, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)

    # Arabic tagline
    add_text(s, Inches(1.5), Inches(3.85), cw - Inches(3), Inches(0.5),
             {"text": "نظام عصبي ذكي للدعم التقني المؤسسي",
              "font": FONT_BODY, "size": 22,
              "italic": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)

    # English tagline
    add_text(s, Inches(2), Inches(4.45), cw - Inches(4), Inches(0.4),
             {"text": "An AI nervous system for Saudi enterprise IT support",
              "font": FONT_BODY, "size": 14, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)

    # Decorative four agent dots with a connecting line
    dot_y = Inches(5.85)
    dot_r = Inches(0.09)
    spacing = Inches(1.0)
    colors = [MEMORY, RESOLVER, GUARDIAN, REPORTER]
    total_w = spacing * (len(colors) - 1)
    start_x = cw // 2 - total_w // 2
    add_hline(s, start_x, dot_y, start_x + total_w, color=DIVIDER, weight=0.75)
    for i, c in enumerate(colors):
        add_dot(s, start_x + spacing * i, dot_y, dot_r, c)

    # Team
    add_text(s, Inches(2), Inches(6.4), cw - Inches(4), Inches(0.4),
             {"text": "Haneen Aldossari    •    Noura Aldossari",
              "font": FONT_BODY, "size": 14, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=False)
    add_text(s, Inches(2), Inches(6.8), cw - Inches(4), Inches(0.3),
             {"text": "April 2026",
              "font": FONT_BODY, "size": 12, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)


def slide_problem(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "المشكلة", "The Problem")

    # Quote (italic, left+right padded)
    add_text(s, Inches(0.7), Inches(2.0), SLIDE_W - Inches(1.4), Inches(1.0),
             {"text":
              "“موظف يفتح تذكرة دعم لمشكلة بسيطة، وينتظر ٣ ساعات قبل أن "
              "تصل الإجابة. هذا واقع يومي في الشركات السعودية.”",
              "font": FONT_BODY, "size": 18,
              "italic": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.RIGHT, rtl=True)

    # Each problem: short Arabic heading → small English subtitle → Arabic body.
    # Tighter than the v1 wording — every line earns its place.
    points = [
        ("١.  مشاكل متكررة تستهلك ساعات الدعم",
         "Repetitive tickets eat IT hours",
         "٨٠٪ من التذاكر اليومية تتكرر — VPN، password، email sync. "
         "كل واحدة تأخذ ساعات، معظمها انتظار في طوابير."),
        ("٢.  فجوة لغوية في الأدوات الحالية",
         "English-only tools miss the Saudi mix",
         "الموظف يتحدث عربي ويستخدم المصطلحات التقنية بالإنجليزية بشكل "
         "طبيعي. الأدوات الموجودة لا تفهم هذا المزج المحلي."),
        ("٣.  معرفة تضيع مع كل تذكرة تُحَل",
         "No collective memory, no NCA audit trail",
         "لا ذاكرة مؤسسية، لا اكتشاف للأنماط، ولا توثيق آلي لمتطلبات "
         "NCA Essential Cybersecurity Controls."),
    ]
    point_top = Inches(3.15)
    point_h = Inches(1.15)
    point_w = SLIDE_W - Inches(1.4)
    for i, (head_ar, sub_en, body) in enumerate(points):
        y = point_top + point_h * i
        add_text(s, Inches(0.7), y, point_w, Inches(0.4),
                 {"text": head_ar, "font": FONT_HEAD, "size": 19,
                  "bold": True, "color": TEXT_PRIMARY},
                 align=PP_ALIGN.RIGHT, rtl=True)
        add_text(s, Inches(0.7), y + Inches(0.42), point_w, Inches(0.28),
                 {"text": sub_en, "font": FONT_BODY, "size": 11,
                  "italic": True, "color": TEXT_SECONDARY},
                 align=PP_ALIGN.RIGHT, rtl=False)
        add_text(s, Inches(0.7), y + Inches(0.72), point_w, Inches(0.42),
                 {"text": body, "font": FONT_BODY, "size": 13,
                  "color": TEXT_SECONDARY},
                 align=PP_ALIGN.RIGHT, rtl=True)

    # Result
    add_text(s, Inches(0.7), Inches(6.65), SLIDE_W - Inches(1.4), Inches(0.5),
             {"text":
              "النتيجة: شركات تنفق الملايين، موظفون محبَطون، وفرق IT غارقة "
              "في المتكرر بدلاً من الهندسة الحقيقية.",
              "font": FONT_BODY, "size": 14,
              "italic": True, "color": TEXT_SECONDARY},
             align=PP_ALIGN.RIGHT, rtl=True)


def _agent_card(slide, x, y, w, h, accent, name_en, name_ar, body):
    """
    Two-line header per agent card — Arabic on top (Arabic-first deck),
    English on the line below in muted gray. No parens, no awkward
    English-Arabic mixing on the same line.
    """
    add_card(slide, x, y, w, h, accent=accent)
    pad = Inches(0.28)

    # Arabic name (dominant, in agent's accent color)
    add_text(slide, x + pad, y + Inches(0.27), w - pad * 2, Inches(0.5),
             {"text": name_ar, "font": FONT_HEAD, "size": 22,
              "bold": True, "color": accent},
             align=PP_ALIGN.RIGHT, rtl=True)
    # English name (secondary, muted)
    add_text(slide, x + pad, y + Inches(0.78), w - pad * 2, Inches(0.32),
             {"text": name_en, "font": FONT_BODY, "size": 13,
              "color": TEXT_SECONDARY},
             align=PP_ALIGN.RIGHT, rtl=False)
    # Body
    add_text(slide, x + pad, y + Inches(1.18), w - pad * 2, h - Inches(1.35),
             {"text": body, "font": FONT_BODY, "size": 13,
              "color": TEXT_PRIMARY},
             align=PP_ALIGN.RIGHT, rtl=True)


def slide_solution(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "الحل", "The Solution")

    add_text(s, Inches(0.7), Inches(1.85), SLIDE_W - Inches(1.4), Inches(0.7),
             {"text":
              "مساعد دعم تقني ذكي يعيش على شاشة الموظف ويتفاعل معه بالعربية "
              "والإنجليزية. يضغط زراً واحداً، يصف مشكلته بصوته، ويحصل على "
              "إرشاد بصري وتوجيه فوري بلغته.",
              "font": FONT_BODY, "size": 16, "color": TEXT_PRIMARY},
             align=PP_ALIGN.RIGHT, rtl=True)

    add_text(s, Inches(0.7), Inches(2.65), SLIDE_W - Inches(1.4), Inches(0.4),
             {"text":
              "خلف الكواليس، أربعة وكلاء يتعاونون كفريق متكامل لحل المشكلة:",
              "font": FONT_BODY, "size": 16, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=True)

    # 2x2 grid
    card_w = (SLIDE_W - Inches(1.6)) // 2 - Inches(0.15)
    card_h = Inches(1.85)
    left_x = Inches(0.7)
    right_x = SLIDE_W - Inches(0.7) - card_w
    top_y = Inches(3.25)
    bot_y = top_y + card_h + Inches(0.22)

    _agent_card(s, right_x, top_y, card_w, card_h, MEMORY,
                "Memory", "وكيل الذاكرة",
                "يستذكر البلاغات السابقة ويكتشف الأنماط في الشركة. يربط "
                "البلاغات المتعددة عبر موظفين مختلفين لاكتشاف الحوادث "
                "الجماعية.")
    _agent_card(s, left_x, top_y, card_w, card_h, RESOLVER,
                "Resolver", "وكيل المُحلِّل",
                "يقرأ الشاشة، يفهم سؤال الموظف، ويقترح خطوة عملية واضحة "
                "مع تحديد الموقع الذي يجب الضغط عليه.")
    _agent_card(s, right_x, bot_y, card_w, card_h, GUARDIAN,
                "Guardian", "وكيل الحارس",
                "يراجع كل إجراء أمام سياسات الشركة وضوابط الـ NCA قبل "
                "تنفيذه. يملك صلاحية رفض الإجراءات الخطرة واقتراح بدائل "
                "آمنة.")
    _agent_card(s, left_x, bot_y, card_w, card_h, REPORTER,
                "Reporter", "وكيل المُبلِّغ",
                "يجمّع الرد النهائي للموظف ويُنبّه الفريق المختص تلقائياً، "
                "ويولّد سجلاً مدقَّقاً لكل عملية يطابق متطلبات NCA.")


def slide_how(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "كيف يعمل النظام", "How It Works")

    stages = [
        # (en_title, ar_title, ar_body, en_body)
        ("User", "المستخدم",
         "يضغط hotkey ويصف المشكلة بصوته",
         "User presses hotkey, speaks Arabic"),
        ("Capture", "التقاط",
         "لقطة شاشة + تحويل صوت لنص",
         "Screenshot + voice-to-text"),
        ("Agents Collaborate", "الوكلاء يتعاونون",
         "Memory → Resolver → Guardian → Reporter", ""),
        ("Response", "الاستجابة",
         "رد عربي + المؤشر يطير على الحل",
         "Arabic response + cursor flies to solution"),
    ]
    n = len(stages)
    margin = Inches(0.55)
    arrow_w = Inches(0.4)
    avail = SLIDE_W - margin * 2
    stage_w = (avail - arrow_w * (n - 1)) / n
    stage_y = Inches(2.1)
    stage_h = Inches(1.85)

    for i, (en, ar, ar_body, en_body) in enumerate(stages):
        # RTL: i=0 (User) on the right, i=3 (Response) on the left
        x = SLIDE_W - margin - stage_w - i * (stage_w + arrow_w)
        add_card(s, x, stage_y, stage_w, stage_h)

        add_text(s, x + Inches(0.1), stage_y + Inches(0.25),
                 stage_w - Inches(0.2), Inches(0.45),
                 {"text": ar, "font": FONT_HEAD, "size": 16,
                  "bold": True, "color": TEXT_PRIMARY},
                 align=PP_ALIGN.CENTER, rtl=True)
        add_text(s, x + Inches(0.1), stage_y + Inches(0.7),
                 stage_w - Inches(0.2), Inches(0.3),
                 {"text": en, "font": FONT_BODY, "size": 11,
                  "color": TEXT_SECONDARY},
                 align=PP_ALIGN.CENTER, rtl=False)
        add_text(s, x + Inches(0.1), stage_y + Inches(1.05),
                 stage_w - Inches(0.2), Inches(0.4),
                 {"text": ar_body, "font": FONT_BODY, "size": 11,
                  "color": TEXT_PRIMARY},
                 align=PP_ALIGN.CENTER, rtl=True)
        if en_body:
            add_text(s, x + Inches(0.1), stage_y + Inches(1.45),
                     stage_w - Inches(0.2), Inches(0.3),
                     {"text": en_body, "font": FONT_BODY, "size": 9,
                      "color": TEXT_SECONDARY},
                     align=PP_ALIGN.CENTER, rtl=False)

        # Stage 3: 4 colored dots representing the agents
        if i == 2:
            cx = x + stage_w / 2
            dy = stage_y + Inches(1.55)
            dot_r = Inches(0.06)
            spacing = Inches(0.18)
            colors = [MEMORY, RESOLVER, GUARDIAN, REPORTER]
            total = spacing * (len(colors) - 1)
            sx = cx - total / 2
            for j, c in enumerate(colors):
                add_dot(s, sx + spacing * j, dy, dot_r, c)

        # Arrow LEFT-pointing between stages (RTL flow)
        if i < n - 1:
            ax = x - arrow_w
            add_text(s, ax, stage_y + Inches(0.75),
                     arrow_w, Inches(0.45),
                     {"text": "←", "font": FONT_BODY, "size": 28,
                      "color": TEXT_SECONDARY},
                     align=PP_ALIGN.CENTER, rtl=False)

    # Section heading — Arabic-first
    add_text(s, Inches(0.7), Inches(4.25), SLIDE_W - Inches(1.4), Inches(0.55),
             {"text":
              "وكلاء يتعاونون كفريق متكامل لحل المشكلة",
              "font": FONT_HEAD, "size": 22,
              "bold": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)
    add_text(s, Inches(0.7), Inches(4.78), SLIDE_W - Inches(1.4), Inches(0.35),
             {"text":
              "Four agents collaborating as one team",
              "font": FONT_BODY, "size": 13,
              "italic": True, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)

    bullets = [
        "الذاكرة قد تخبر المُحلِّل بتجاوز التشخيص لأن المشكلة معروفة",
        "الحارس قد يرفض اقتراح المُحلِّل ويرغمه على البحث عن بديل آمن",
        "المُبلِّغ قد يربط الطلب بحادثة جماعية قائمة بدلاً من إنشاء "
        "تذكرة جديدة",
    ]
    by = Inches(5.4)
    bh = Inches(0.5)
    for i, txt in enumerate(bullets):
        add_text(s, Inches(0.7), by + bh * i,
                 SLIDE_W - Inches(1.4), Inches(0.5),
                 {"text": "●  ", "font": FONT_BODY, "size": 14,
                  "color": TEXT_SECONDARY},
                 {"text": txt, "font": FONT_BODY, "size": 16,
                  "color": TEXT_PRIMARY},
                 align=PP_ALIGN.RIGHT, rtl=True)


def slide_diff(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "لماذا نحن مختلفون", "Why We Win")

    cards = [
        ("Truly Multi-Agent",
         "تعاون حقيقي بين الوكلاء",
         "أربعة وكلاء يتعاونون كفريق متكامل لحل المشكلة. الذاكرة "
         "تخبر المُحلِّل بما هو معروف، الحارس يعدّل الخطة عند الحاجة، "
         "والمُبلِّغ يعيد توجيه المسار بالكامل عند اكتشاف نمط جماعي."),
        ("Arabic-First Design",
         "مصمَّم للعربية أولاً",
         "تصميم مبني من الأساس للسياق السعودي. واجهة RTL كاملة، وفهم "
         "للمصطلحات التقنية الإنجليزية كما يستخدمها الموظف السعودي "
         "طبيعياً في كلامه اليومي."),
        ("NCA-Compliant by Default",
         "متوافق مع NCA افتراضياً",
         "سجل مدقَّق لكل عملية يطابق Essential Cybersecurity Controls — "
         "جزء أساسي من تصميم الـ Reporter agent منذ اليوم الأول، لا "
         "ميزة لاحقة."),
        ("Open-Source Foundation",
         "أساس مفتوح المصدر",
         "نظام شفاف بالكامل، لا قفل على بائع، وقابلية تطوير محلية "
         "للسوق السعودي."),
    ]

    card_w = (SLIDE_W - Inches(1.6)) / 2 - Inches(0.15)
    card_h = Inches(1.85)
    left_x = Inches(0.7)
    right_x = SLIDE_W - Inches(0.7) - card_w
    top_y = Inches(2.0)
    bot_y = top_y + card_h + Inches(0.22)
    positions = [(right_x, top_y), (left_x, top_y),
                 (right_x, bot_y), (left_x, bot_y)]

    for (en, ar, body), (x, y) in zip(cards, positions):
        add_card(s, x, y, card_w, card_h)
        add_text(s, x + Inches(0.28), y + Inches(0.27),
                 card_w - Inches(0.56), Inches(0.5),
                 {"text": en, "font": FONT_HEAD, "size": 18,
                  "bold": True, "color": TEXT_PRIMARY},
                 align=PP_ALIGN.RIGHT, rtl=False)
        add_text(s, x + Inches(0.28), y + Inches(0.7),
                 card_w - Inches(0.56), Inches(0.4),
                 {"text": ar, "font": FONT_BODY, "size": 14,
                  "italic": True, "color": TEXT_SECONDARY},
                 align=PP_ALIGN.RIGHT, rtl=True)
        add_text(s, x + Inches(0.28), y + Inches(1.1),
                 card_w - Inches(0.56), card_h - Inches(1.25),
                 {"text": body, "font": FONT_BODY, "size": 12,
                  "color": TEXT_PRIMARY},
                 align=PP_ALIGN.RIGHT, rtl=True)

    add_text(s, Inches(0.7), Inches(6.55), SLIDE_W - Inches(1.4), Inches(0.6),
             {"text":
              "“نظام دعم تقني عربي، يفهم شركتك، ويتعلم من كل بلاغ يحلّه.”",
              "font": FONT_BODY, "size": 18,
              "italic": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)


def slide_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "خارطة الطريق", "Roadmap")

    # Three columns: Before (done) → During (in progress) → After (planned).
    # Each gets a status badge under the bilingual header so the phase
    # state reads at a glance.
    phases = [
        {
            "ar": "قبل الهاكاثون",
            "en": "Before",
            "badge": "✅  مكتمل",
            "accent": RESOLVER,  # green
            "items": [
                "تحليل السوق وتصميم البنية متعددة الوكلاء",
                "توثيق تقني كامل (PRD + Agent Specs)",
                "Fork لـ Flicky وتعريب الواجهة",
                "بناء Agent Panel البصري",
            ],
        },
        {
            "ar": "خلال الهاكاثون",
            "en": "During (3 Days)",
            "badge": "⏳  قيد التنفيذ",
            "accent": MEMORY,  # blue
            "items": [
                "تفعيل الوكلاء الأربعة بـ Claude API",
                "بناء طبقة RAG لقاعدة معرفة الشركة",
                "سيناريو اكتشاف الحوادث الجماعية",
                "سيناريو مراجعة Guardian",
                "ديمو شغّال end-to-end",
            ],
        },
        {
            "ar": "ما بعد الهاكاثون",
            "en": "After",
            "badge": "🚀  مخطط",
            "accent": TEXT_SECONDARY,  # gray
            "items": [
                "Pilot مع شركة سعودية",
                "تكامل مع Slack و Teams",
                "لوحة تقارير NCA compliance",
                "نسخة SaaS متعددة المستأجرين",
            ],
        },
    ]
    margin = Inches(0.6)
    spacing = Inches(0.25)
    col_w = (SLIDE_W - margin * 2 - spacing * 2) / 3
    col_y = Inches(1.85)
    col_h = Inches(5.2)

    for i, phase in enumerate(phases):
        # RTL: i=0 (Before) on the right, i=2 (After) on the left.
        x = SLIDE_W - margin - col_w - i * (col_w + spacing)
        add_card(s, x, col_y, col_w, col_h, accent=phase["accent"])
        pad = Inches(0.22)

        # Bilingual header — Arabic dominant, English secondary
        add_text(s, x + pad, col_y + Inches(0.22),
                 col_w - pad * 2, Inches(0.45),
                 {"text": phase["ar"],
                  "font": FONT_HEAD, "size": 17,
                  "bold": True, "color": phase["accent"]},
                 align=PP_ALIGN.RIGHT, rtl=True)
        add_text(s, x + pad, col_y + Inches(0.7),
                 col_w - pad * 2, Inches(0.3),
                 {"text": phase["en"],
                  "font": FONT_BODY, "size": 11, "color": TEXT_SECONDARY},
                 align=PP_ALIGN.RIGHT, rtl=False)

        # Status badge — small accent-colored pill of text
        add_text(s, x + pad, col_y + Inches(1.05),
                 col_w - pad * 2, Inches(0.35),
                 {"text": phase["badge"],
                  "font": FONT_BODY, "size": 12, "bold": True,
                  "color": phase["accent"]},
                 align=PP_ALIGN.RIGHT, rtl=True)

        # Items
        item_y = col_y + Inches(1.6)
        item_h = Inches(0.6)
        for j, item in enumerate(phase["items"]):
            add_text(s, x + pad, item_y + item_h * j,
                     col_w - pad * 2, item_h,
                     {"text": "●  ", "font": FONT_BODY, "size": 11,
                      "color": phase["accent"]},
                     {"text": item, "font": FONT_BODY, "size": 12,
                      "color": TEXT_PRIMARY},
                     align=PP_ALIGN.RIGHT, rtl=True)


def slide_team(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_solid_bg(s)
    add_slide_title(s, "الفريق", "The Team")

    members = [
        ("Haneen Aldossari", "CS student at PSAU"),
        ("Noura Aldossari",  "CS student at PSAU"),
    ]

    card_w = (SLIDE_W - Inches(1.6)) / 2 - Inches(0.2)
    card_h = Inches(1.85)
    left_x = Inches(0.7)
    right_x = SLIDE_W - Inches(0.7) - card_w
    card_y = Inches(2.3)

    for (name, line), x in zip(members, [right_x, left_x]):
        add_card(s, x, card_y, card_w, card_h)
        add_text(s, x + Inches(0.3), card_y + Inches(0.55),
                 card_w - Inches(0.6), Inches(0.6),
                 {"text": name, "font": FONT_HEAD, "size": 24,
                  "bold": True, "color": TEXT_PRIMARY},
                 align=PP_ALIGN.CENTER, rtl=False)
        add_text(s, x + Inches(0.3), card_y + Inches(1.15),
                 card_w - Inches(0.6), Inches(0.4),
                 {"text": line, "font": FONT_BODY, "size": 14,
                  "color": TEXT_SECONDARY},
                 align=PP_ALIGN.CENTER, rtl=False)

    # Closing tagline (Built On section removed)
    add_text(s, Inches(0.7), Inches(5.4),
             SLIDE_W - Inches(1.4), Inches(0.7),
             {"text": "نظام دعم تقني عربي، مصمَّم للسوق السعودي",
              "font": FONT_HEAD, "size": 28,
              "bold": True, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)
    add_text(s, Inches(0.7), Inches(6.15),
             SLIDE_W - Inches(1.4), Inches(0.45),
             {"text": "متوافق مع رؤية ٢٠٣٠",
              "font": FONT_BODY, "size": 20, "color": TEXT_PRIMARY},
             align=PP_ALIGN.CENTER, rtl=True)
    add_text(s, Inches(0.7), Inches(6.7),
             SLIDE_W - Inches(1.4), Inches(0.4),
             {"text": "Saudi-first IT support, built for Vision 2030",
              "font": FONT_BODY, "size": 13,
              "italic": True, "color": TEXT_SECONDARY},
             align=PP_ALIGN.CENTER, rtl=False)


# ── Main ─────────────────────────────────────────────────────────────────


def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_how(prs)
    slide_diff(prs)
    slide_roadmap(prs)
    slide_team(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    size = OUT.stat().st_size
    print(f"Wrote {OUT}")
    print(f"Size: {size:,} bytes ({size / 1024:.1f} KB)")


if __name__ == "__main__":
    build()
